import os
import re
import unicodedata
from pptx import Presentation
from core.ports.secondary.services import PptxFileReadingService
from core.entities import InputFile, FileContent, PageContent, PptxReadFileConfig

class PptxFileReadingServiceImpl(PptxFileReadingService):
    """Service implementation for reading PPTX files."""
    def __init__(self):
        pass

    def read_file(self, input_file: InputFile, read_file_config: dict, **kwargs) -> FileContent:
        ppt = Presentation(input_file.local_file_path)
        docs = []
        for i, slide in enumerate(ppt.slides):
            content = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text
            refined_content = unicodedata.normalize("NFKC", self.remove_extra_space(content))
            docs.append(refined_content)
        contents = [f"# Slide {i+1}: {doc}" for i, doc in enumerate(docs)]
        page_content = PageContent(
            content="\n\n".join(contents),
            page_number=0
        )
        page_contents = [page_content]
        file_content = FileContent(
            file_name=input_file.file_name,
            file_path=input_file.local_file_path,
            page_contents=page_contents,
            content_format="markdown",
        )
        return file_content

    def remove_extra_space(self, s):
        pattern = r"(?<!\S) +| +(?!\S)"
        s = re.sub(pattern, " ", s)
        return s

    async def aread_file(self, input_file, read_file_config, **kwargs):
        # TODO: Implement asynchronous reading of PPTX files
        pass
